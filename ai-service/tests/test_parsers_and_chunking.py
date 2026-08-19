"""
Document parsing and structure-aware chunking.

Two things are being protected here. First, that structure survives parsing —
a heading, sheet name or slide number is what makes a citation say "Budget.xlsx,
sheet Q3" instead of just the filename, and it cannot be recovered downstream if
lost. Second, that malformed input fails as a readable error rather than an
unhandled exception, since these parsers read attacker-supplied binaries.
"""
from __future__ import annotations

import io

import pytest

from rag import parsers
from rag.chunking import chunk_blocks


# ── Word ────────────────────────────────────────────────────────────────────

@pytest.fixture()
def docx_bytes() -> bytes:
    import docx

    document = docx.Document()
    document.add_heading("Leave Policy", level=1)
    document.add_paragraph("Employees receive 25 days of annual leave per year.")
    document.add_heading("Remote Work", level=1)
    document.add_paragraph("Remote work is permitted three days per week.")

    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Region"
    table.cell(0, 1).text = "Allowance"
    table.cell(1, 0).text = "EMEA"
    table.cell(1, 1).text = "28 days"

    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def test_docx_paragraphs_are_parsed(docx_bytes):
    text = " ".join(b.text for b in parsers.parse("word", docx_bytes, "policy.docx"))
    assert "25 days of annual leave" in text


def test_docx_headings_become_metadata(docx_bytes):
    blocks = parsers.parse("word", docx_bytes, "policy.docx")
    assert any(b.metadata.get("heading") == "Remote Work" for b in blocks)


def test_docx_tables_are_flattened_row_wise(docx_bytes):
    text = " ".join(b.text for b in parsers.parse("word", docx_bytes, "policy.docx"))
    assert "EMEA | 28 days" in text


# ── Excel ───────────────────────────────────────────────────────────────────

@pytest.fixture()
def xlsx_bytes() -> bytes:
    from openpyxl import Workbook

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Q3"
    sheet.append(["Region", "Revenue", "Owner"])
    sheet.append(["EMEA", "1.2M", "Priya"])
    sheet.append(["APAC", "0.9M", "Chen"])

    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def test_xlsx_rows_are_self_describing(xlsx_bytes):
    """
    A retrieved row has to survive out of context. "EMEA | 1.2M" is meaningless
    on its own; "Region: EMEA | Revenue: 1.2M" is not.
    """
    text = " ".join(b.text for b in parsers.parse("excel", xlsx_bytes, "budget.xlsx"))
    assert "Region: EMEA" in text
    assert "Revenue: 1.2M" in text


def test_xlsx_sheet_name_is_kept(xlsx_bytes):
    blocks = parsers.parse("excel", xlsx_bytes, "budget.xlsx")
    assert any(b.metadata.get("sheet") == "Q3" for b in blocks)


# ── PowerPoint ──────────────────────────────────────────────────────────────

@pytest.fixture()
def pptx_bytes() -> bytes:
    from pptx import Presentation

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Roadmap 2026"
    slide.placeholders[1].text = "Ship billing migration in Q1"

    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def test_pptx_slide_text_is_parsed(pptx_bytes):
    text = " ".join(b.text for b in parsers.parse("powerpoint", pptx_bytes, "deck.pptx"))
    assert "billing migration" in text


def test_pptx_keeps_slide_number_and_title(pptx_bytes):
    blocks = parsers.parse("powerpoint", pptx_bytes, "deck.pptx")
    assert any(b.metadata.get("page") == 1 and b.metadata.get("heading") == "Roadmap 2026"
               for b in blocks)


# ── Images ──────────────────────────────────────────────────────────────────

def test_image_is_indexed_with_an_honest_caveat():
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (120, 60), "navy").save(buffer, format="PNG")

    blocks = parsers.parse("image", buffer.getvalue(), "chart.png")
    text = blocks[0].text
    assert "120x60" in text
    # OCR is not implemented; saying so is better than silently indexing
    # nothing and letting the model imply it read the image.
    assert "not been transcribed" in text


# ── Failure handling ────────────────────────────────────────────────────────

def test_corrupt_office_file_raises_a_handled_error():
    with pytest.raises(Exception):
        parsers.parse("word", b"not really a docx at all", "fake.docx")


def test_empty_file_is_rejected():
    with pytest.raises(parsers.UnsupportedFile):
        parsers.parse("text", b"   ", "empty.txt")


def test_unknown_type_is_rejected():
    with pytest.raises(parsers.UnsupportedFile):
        parsers.parse("executable", b"MZ", "x.exe")


# ── Chunking ────────────────────────────────────────────────────────────────

def test_small_blocks_in_the_same_context_are_merged():
    """
    One spreadsheet row is far below the token budget. Embedding hundreds of
    one-line chunks produces uniformly weak matches, so adjacent rows merge.
    """
    blocks = [
        parsers.Block("Row A detail", {"sheet": "Q3", "row": 2}),
        parsers.Block("Row B detail", {"sheet": "Q3", "row": 3}),
    ]
    chunks = chunk_blocks(blocks, source_path="budget.xlsx")
    assert any("Row A detail" in c.text and "Row B detail" in c.text for c in chunks)


def test_blocks_from_different_contexts_are_never_merged():
    """A chunk straddling two sheets would cite the wrong one."""
    blocks = [
        parsers.Block("Row A detail", {"sheet": "Q3", "row": 2}),
        parsers.Block("Other sheet", {"sheet": "Q4", "row": 2}),
    ]
    chunks = chunk_blocks(blocks, source_path="budget.xlsx")
    assert all(not ("Row A detail" in c.text and "Other sheet" in c.text) for c in chunks)


def test_chunks_carry_source_and_structure():
    blocks = [parsers.Block("Some content", {"sheet": "Q4", "row": 2})]
    chunks = chunk_blocks(blocks, source_path="budget.xlsx")
    assert chunks
    assert all(c.metadata["source_path"] == "budget.xlsx" for c in chunks)
    assert any(c.metadata.get("sheet") == "Q4" for c in chunks)


def test_oversized_block_is_split_rather_than_dropped():
    huge = parsers.Block("word " * 4000, {"page": 1})
    chunks = chunk_blocks([huge], source_path="big.pdf", max_tokens=200)
    assert len(chunks) > 1
    assert all(c.metadata.get("page") == 1 for c in chunks)
