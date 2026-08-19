"""
rag/parsers.py — turn an uploaded file into structured text blocks.

Each parser returns a list of `Block`s rather than one string, because the
chunker needs the structure to do its job: a heading, a slide number or a sheet
name is what makes a later citation say "Budget.xlsx, sheet Q3" instead of
"Budget.xlsx". Losing that at parse time cannot be recovered downstream.

Every parser is defensive. These libraries read attacker-supplied binaries, so
a malformed file must produce an error we report, never an unhandled exception
that takes the worker down. Optional imports are lazy so a missing extra
disables one format instead of breaking startup.
"""
from __future__ import annotations

import io
import logging
from dataclasses import dataclass, field
from typing import Any

log = logging.getLogger(__name__)

# Guards against a decompression bomb: a 2 MB .xlsx can expand into millions of
# cells, and a 50 MB .pptx into thousands of slides.
MAX_BLOCKS = 5_000
MAX_CHARS = 4_000_000
MAX_SHEET_ROWS = 5_000
MAX_SHEET_COLS = 100


@dataclass
class Block:
    """One structural unit of a document."""
    text: str
    metadata: dict[str, Any] = field(default_factory=dict)


class UnsupportedFile(ValueError):
    """The file type is not handled, or its optional dependency is absent."""


def _cap(blocks: list[Block]) -> list[Block]:
    total = 0
    kept: list[Block] = []
    for block in blocks[:MAX_BLOCKS]:
        total += len(block.text)
        if total > MAX_CHARS:
            log.warning("Document truncated at %d characters.", MAX_CHARS)
            break
        kept.append(block)
    return kept


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------

def parse_pdf(data: bytes) -> list[Block]:
    import fitz  # PyMuPDF

    blocks: list[Block] = []
    with fitz.open(stream=data, filetype="pdf") as doc:
        for number, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                blocks.append(Block(text, {"page": number}))
    return _cap(blocks)


# ---------------------------------------------------------------------------
# Plain text / Markdown
# ---------------------------------------------------------------------------

def parse_text(data: bytes) -> list[Block]:
    return _cap([Block(data.decode("utf-8", errors="replace"))])


# ---------------------------------------------------------------------------
# Word
# ---------------------------------------------------------------------------

def parse_docx(data: bytes) -> list[Block]:
    try:
        import docx  # python-docx
    except ImportError as e:  # pragma: no cover
        raise UnsupportedFile("Word support needs `python-docx`.") from e

    document = docx.Document(io.BytesIO(data))
    blocks: list[Block] = []
    heading = None

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        # Word marks headings by style, which is a far better section signal
        # than guessing from font size or position.
        if (paragraph.style.name or "").lower().startswith("heading"):
            heading = text
            continue
        blocks.append(Block(text, {"heading": heading}))

    # Tables carry a lot of a typical business document's actual content;
    # flattening them row-wise keeps each row independently retrievable.
    for index, table in enumerate(document.tables, start=1):
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                blocks.append(Block(" | ".join(cells), {"heading": heading, "table": index}))

    return _cap(blocks)


# ---------------------------------------------------------------------------
# Excel
# ---------------------------------------------------------------------------

def parse_xlsx(data: bytes) -> list[Block]:
    try:
        from openpyxl import load_workbook
    except ImportError as e:  # pragma: no cover
        raise UnsupportedFile("Excel support needs `openpyxl`.") from e

    # read_only + data_only: stream rows instead of building the whole object
    # model, and take computed values rather than formula text.
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    blocks: list[Block] = []

    try:
        for sheet in workbook.worksheets:
            header: list[str] = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= MAX_SHEET_ROWS:
                    break
                values = ["" if v is None else str(v).strip() for v in row[:MAX_SHEET_COLS]]
                if not any(values):
                    continue

                if not header:
                    header = values
                    continue

                # Pair each value with its column header so a retrieved row is
                # self-describing: "Region: EMEA | Revenue: 1.2M" survives out
                # of context, where "EMEA | 1.2M" does not.
                pairs = [
                    f"{head}: {value}"
                    for head, value in zip(header, values)
                    if value and head
                ]
                if pairs:
                    blocks.append(Block(" | ".join(pairs), {"sheet": sheet.title, "row": index + 1}))
    finally:
        workbook.close()

    return _cap(blocks)


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def parse_pptx(data: bytes) -> list[Block]:
    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover
        raise UnsupportedFile("PowerPoint support needs `python-pptx`.") from e

    presentation = Presentation(io.BytesIO(data))
    blocks: list[Block] = []

    for number, slide in enumerate(presentation.slides, start=1):
        pieces: list[str] = []
        title = None

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None:
                title = text.splitlines()[0][:120]
            pieces.append(text)

        if pieces:
            # One block per slide: a slide is the natural unit a person cites.
            blocks.append(Block("\n".join(pieces), {"page": number, "heading": title}))

    return _cap(blocks)


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

def parse_image(data: bytes, filename: str = "") -> list[Block]:
    """
    Images carry no extractable text without OCR, which is deliberately out of
    scope. Rather than failing the upload, index what is genuinely known — the
    filename and the image's own properties — so the file is at least
    discoverable, and say plainly that the contents were not read.
    """
    try:
        from PIL import Image
    except ImportError as e:  # pragma: no cover
        raise UnsupportedFile("Image support needs `pillow`.") from e

    with Image.open(io.BytesIO(data)) as image:
        image.verify()  # structural check; raises on a malformed file
        width, height = image.size
        fmt = (image.format or "image").upper()

    label = filename or "image"
    return [Block(
        f"Image file: {label}. Format {fmt}, {width}x{height} pixels. "
        f"The visual contents of this image have not been transcribed — OCR is "
        f"not enabled, so questions about text inside this image cannot be "
        f"answered from it.",
        {"heading": label, "kind": "image"},
    )]


# ---------------------------------------------------------------------------
# Dispatch
# ---------------------------------------------------------------------------

PARSERS = {
    "pdf": parse_pdf,
    "markdown": parse_text,
    "text": parse_text,
    "word": parse_docx,
    "excel": parse_xlsx,
    "powerpoint": parse_pptx,
}


def parse(file_type: str, data: bytes, filename: str = "") -> list[Block]:
    """Parse `data` according to its logical type. Raises UnsupportedFile."""
    if file_type == "image":
        return parse_image(data, filename)

    parser = PARSERS.get(file_type)
    if parser is None:
        raise UnsupportedFile(f"Unsupported file type: {file_type}")

    blocks = parser(data)
    if not blocks or not any(b.text.strip() for b in blocks):
        raise UnsupportedFile(
            "No readable text was found in this file. If it is a scan or a "
            "photo, OCR is not supported yet."
        )
    return blocks
